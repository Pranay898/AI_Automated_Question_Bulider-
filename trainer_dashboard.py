import docx2txt
from pptx import Presentation
import pytesseract
import streamlit as st
from db import curricula_collection, generated_questions_collection
from datetime import datetime
from utils import generate_questions
import os
import csv
from io import StringIO, BytesIO
from fpdf import FPDF
import fitz  # PyMuPDF for PDF processing; install with 'pip install pymupdf'
from PIL import Image
import pandas as pd
from ques_bank_gen import gen_que

api_key = "AIzaSyCkhwdytqyr039-BnhACY2RzexSUZZcwB4"
os.environ["GOOGLE_API_KEY"] = api_key

def extract_text_from_pdf(file):
    with fitz.open(stream=file, filetype="pdf") as pdf_document:
        text = ""
        for page in pdf_document:
            text += page.get_text()
    return text

def extract_text_from_docx(file):
    return docx2txt.process(file)

def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_image(file):
    image = Image.open(file)
    return pytesseract.image_to_string(image)




def trainer_dashboard():
    st.sidebar.title("Trainer Dashboard")
    option = st.sidebar.selectbox("Trainer Options", ["Upload Curriculum", "Generate Question Bank", "Review and Edit", "Download Question Bank"])

    if option == "Upload Curriculum":
        st.title("Upload Curriculum")
        curriculum_file = st.file_uploader("Upload Curriculum (Any format)", type=["pdf", "jpg", "png", "docx", "pptx", "xlsx", "csv"])
        if st.button("Upload"):
            if curriculum_file is not None:
                file_type = curriculum_file.name.split('.')[-1].lower()
                file_content = None

                # Extract text based on the file type
                if file_type == "pdf":
                    file_content = extract_text_from_pdf(curriculum_file.read())
                elif file_type == "docx":
                    file_content = extract_text_from_docx(curriculum_file)
                elif file_type == "pptx":
                    file_content = extract_text_from_pptx(curriculum_file)
                elif file_type in ["jpg", "jpeg", "png"]:
                    file_content = extract_text_from_image(curriculum_file)
                else:
                    st.error("Unsupported file format. Please upload a valid curriculum file.")

                # Store the extracted content in the database
                if file_content:
                    curricula_collection.add({"content": file_content, "filename": curriculum_file.name})
                    st.success("Curriculum uploaded and content stored successfully!")

                    # Verify if the content is stored correctly
                    stored_file = curricula_collection.where("filename", "==", curriculum_file.name).get()
                    if stored_file:
                        st.success("File content stored in the database successfully!")
                    else:
                        st.error("Failed to store the file content in the database.")
            else:
                st.error("Please upload a valid file.")

    elif option == "Generate Question Bank":
        gen_que()

    elif option == "Review and Edit":
        st.title("Review and Edit Questions")
        existing_question_banks = generated_questions_collection.stream()
        question_bank_ids = [q.id for q in existing_question_banks]

        selected_question_bank = st.selectbox("Select Question Bank to Edit", question_bank_ids)

        if selected_question_bank:
            question_bank_details = generated_questions_collection.document(selected_question_bank).get().to_dict()
            if question_bank_details and "questions" in question_bank_details:
                questions = question_bank_details["questions"]
                question_options = [q["question"] for q in questions]

                selected_question = st.selectbox("Select Question to Edit", question_options)

                if selected_question:
                    current_question = next(q for q in questions if q["question"] == selected_question)
                    
                    # Edit question, answer, and options
                    new_question = st.text_input("Edit Question", value=current_question["question"])
                    new_answer = st.text_input("Edit Answer", value=current_question.get("answer", ""))
                    new_option_1 = st.text_input("Edit Option A)", value=current_question.get("option-1", ""))
                    new_option_2 = st.text_input("Edit Option B)", value=current_question.get("option-2", ""))
                    new_option_3 = st.text_input("Edit Option C)", value=current_question.get("option-3", ""))
                    new_option_4 = st.text_input("Edit Option D)", value=current_question.get("option-4", ""))

                    if st.button("Save Changes"):
                        # Update the question with new values
                        updated_question = {
                            "question": new_question,
                            "answer": new_answer,
                            "option-1": new_option_1,
                            "option-2": new_option_2,
                            "option-3": new_option_3,
                            "option-4": new_option_4
                        }
                        index = questions.index(current_question)
                        questions[index] = updated_question
                        
                        # Update in database
                        generated_questions_collection.document(selected_question_bank).update({"questions": questions})
                        st.success("Question updated successfully!")

    elif option == "Download Question Bank":
        st.title("Download Question Bank")
        question_banks = generated_questions_collection.stream()
        question_bank_ids = [q.id for q in question_banks]

        selected_question_bank = st.selectbox("Select Question Bank to Download", question_bank_ids)

        if selected_question_bank:
            if st.button("Download Selected Question Bank"):
                question_bank_details = generated_questions_collection.document(selected_question_bank).get().to_dict()
                if question_bank_details and "questions" in question_bank_details:
                    questions = question_bank_details["questions"]

                    # Download in original clean text format
                    text_data = "\n".join([
                        f"Q: {q['question']}\nA: {q['answer']}\nOptions:\nA) {q.get('option-1', '')}\nB) {q.get('option-2', '')}\nC) {q.get('option-3', '')}\nD) {q.get('option-4', '')}\n"
                        for q in questions
                    ])
                    st.download_button(
                        label="Download as Text",
                        data=text_data,
                        file_name=f"{selected_question_bank}.txt",
                        mime="text/plain"
                    )
                else:
                    st.warning("No questions available in this question bank.")
